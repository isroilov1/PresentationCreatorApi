import asyncio
from msilib import Directory
import os
import uuid

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Pt
from pptx.dml.color import RGBColor

from TextPage import add_page_slide_only_text
from TitleImageTextPage import add_page_slide_with_title

async def create(presentationEntity):
    black_templates = [6, 7, 8]
    presentation = Presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)
    # Slide 1
    slide1_register = presentation.slide_layouts[0]
    slide1 = presentation.slides.add_slide(slide1_register)

    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = presentationEntity.Theme
    subtitle1.text = presentationEntity.Author

    # Fon rasmini qo'shish
    bg_img_path = f"wwwroot/assets/templates/{presentationEntity.Template}/1.png"
    left_inch = Inches((presentation.slide_width - Inches(16)) / 2)
    top_inch = Inches((presentation.slide_height - Inches(9)) / 2)
    width_inch = Inches(16)
    height_inch = Inches(9)
    slide1.shapes.add_picture(bg_img_path, left_inch, top_inch, width_inch, height_inch)

    size = 28
    if len(presentationEntity.Theme) < 20:
        size = 38
    text_box = slide1.shapes.add_textbox(Inches(3), Inches(2.5), Inches(10), Inches(4))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text_anchor = "middle"  # Matn textbox markazga joylashadi
    text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Matnning rasmning markaziga moslashtirilishi
    p = text_frame.add_paragraph()
    p.text = presentationEntity.Theme
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.italic = False
    p.font.name = 'Amasis MT Pro Black'
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # Matn markazga moslashtirilishi

    # Author
    text_box = slide1.shapes.add_textbox(Inches(7.5), Inches(5.8), Inches(6), Inches(1))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text_anchor = "top"  # Matn textbox markazga joylashadi
    text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Matnning rasmning markaziga moslashtirilishi
    p = text_frame.add_paragraph()
    p.text = presentationEntity.Author
    p.font.size = Pt(25)
    p.font.bold = False
    p.font.italic = True
    p.font.name = 'Amasis MT Pro Black'
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT  # Matn markazga moslashtirilishi

    if presentationEntity.Template in black_templates:
        # Slide ichidagi barcha shapeslarga murojat qilish
        for shape in slide1.shapes:
            if shape.has_text_frame:  # Faqat matn shakllariga murojat qiling
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)  # Oq rang


    # Page 2
    slide1_register = presentation.slide_layouts[0]
    slide1 = presentation.slides.add_slide(slide1_register)

    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = presentationEntity.Theme
    subtitle1.text = presentationEntity.Author

    # Fon rasmini qo'shish
    bg_img_path = f"wwwroot/assets/templates/{presentationEntity.Template}/2.png"
    left_inch = Inches((presentation.slide_width - Inches(16)) / 2)
    top_inch = Inches((presentation.slide_height - Inches(9)) / 2)
    width_inch = Inches(16)
    height_inch = Inches(9)
    slide1.shapes.add_picture(bg_img_path, left_inch, top_inch, width_inch, height_inch)

    text_box = slide1.shapes.add_textbox(Inches(3), Inches(1), Inches(8), Inches(4))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text_anchor = "middle"  # Matn textbox markazga joylashadi
    text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Matnning rasmning markaziga moslashtirilishi
    p = text_frame.add_paragraph()
    lang = str(presentationEntity.Language).lower()
    reja_text = "Reja:"
    if 'eng' in lang:
        reja_text = "Plan:"
    elif 'rus' in lang:
        reja_text = "????:"
    elif 'ger' in lang:
        reja_text = "Planen:"
    elif 'fre' in lang:
        reja_text = "Plan:"
    elif 'kor' in lang:
        reja_text = "??:"
    p.text = reja_text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.italic = False
    p.font.name = 'Amasis MT Pro Black'
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # Matn markazga moslashtirilishi

    # Rejalar
    top = 2.5
    for i in range(3):
        text_box = slide1.shapes.add_textbox(Inches(1.5), Inches(top), Inches(10), Inches(1))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text_anchor = "top"  # Matn textbox markazga joylashadi
        text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Matnning rasmning markaziga moslashtirilishi
        p = text_frame.add_paragraph()
        p.text = presentationEntity.Titles[i]
        p.font.size = Pt(24)
        p.font.bold = False
        p.font.italic = True
        p.font.name = 'Amasis MT Pro Black'
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT  # Matn markazga moslashtirilishi
        top += 1

    if presentationEntity.Template in black_templates:
        # Slide ichidagi barcha shapeslarga murojat qilish
        for shape in slide1.shapes:
            if shape.has_text_frame:  # Faqat matn shakllariga murojat qiling
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)  # Oq rang

    pageNumber = 2
    
    while pageNumber != presentationEntity.PageCount:
        pageNumber += 1
        await add_page_slide_with_title(presentation, presentationEntity, black_templates, pageNumber)
        if pageNumber == presentationEntity.PageCount:
            break

        pageNumber += 1
        await add_page_slide_only_text(presentation, presentationEntity, black_templates, pageNumber)
        if pageNumber == presentationEntity.PageCount:
            break

        pageNumber += 1
        await add_page_slide_only_text(presentation, presentationEntity, black_templates, pageNumber)
        if pageNumber == presentationEntity.PageCount:
            break

    # End page
    slide1_register = presentation.slide_layouts[0]
    slide1 = presentation.slides.add_slide(slide1_register)

    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = presentationEntity.Theme
    subtitle1.text = presentationEntity.Author

    # Fon rasmini qo'shish
    bg_img_path = f"wwwroot/assets/templates/{presentationEntity.Template}/1.png"
    left_inch = Inches((presentation.slide_width - Inches(16)) / 2)
    top_inch = Inches((presentation.slide_height - Inches(9)) / 2)
    width_inch = Inches(16)
    height_inch = Inches(9)
    slide1.shapes.add_picture(bg_img_path, left_inch, top_inch, width_inch, height_inch)

    #  Theme
    size = 38
    text_box = slide1.shapes.add_textbox(Inches(3), Inches(2.5), Inches(10), Inches(4))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text_anchor = "middle"  # Matn textbox markazga joylashadi
    text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Matnning rasmning markaziga moslashtirilishi
    p = text_frame.add_paragraph()
    finish_text = "E'tiboringiz uchun raxmat!"
    if 'eng' in lang:
        finish_text = "Thank you for your attention!"
    elif 'rus' in lang:
        finish_text = "Спасибо за внимание!"
    elif 'ger' in lang:
        finish_text = "Danke für Ihre Aufmerksamkeit!"
    elif 'fre' in lang:
        finish_text = "Merci pour votre attention!"
    elif 'fre' in lang:
        finish_text = "관심을 가져주셔서 감사합니다!"
    p.text = finish_text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.italic = False
    p.font.name = 'Trade Gothic Next Heavy'
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # Matn markazga moslashtirilishi

    if presentationEntity.Template in black_templates:
        # Slide ichidagi barcha shapeslarga murojat qilish
        for shape in slide1.shapes:
            if shape.has_text_frame:  # Faqat matn shakllariga murojat qiling
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)  # Oq rang

    path_save = f"wwwroot/uploads/presentation/Taqdimot-{uuid.uuid4()}.pptx"
    presentation.save(path_save)
    return path_save


def run_async_function(presentationEntity):
    loop = asyncio.get_event_loop()
    return loop.run_until_complete(create(presentationEntity))
