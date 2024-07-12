
from asyncio.windows_events import INFINITE


async def add_page_slide_only_text(presentation, presentationEntity, black_templates, pageNumber):
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE
    from pptx.util import Inches, Pt
    from PIL import Image
    from pptx.dml.color import RGBColor

    slide1_register = presentation.slide_layouts[0]
    slide1 = presentation.slides.add_slide(slide1_register)

    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = '...'
    subtitle1.text = 't.me/slaydai_bot'

    # Fon rasmini qo'shish
    bg_img_path = f"wwwroot/assets/templates/{presentationEntity.Template}/4.png"
    left_inch = Inches((presentation.slide_width - Inches(16)) / 2)
    top_inch = Inches((presentation.slide_height - Inches(9)) / 2)
    width_inch = Inches(16)
    height_inch = Inches(9)
    slide1.shapes.add_picture(bg_img_path, left_inch, top_inch, width_inch, height_inch)

    # Matn
    text_box = slide1.shapes.add_textbox(Inches(3.5), Inches(1.7), Inches(9), Inches(5.8))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text_anchor = "top"  # Matn textbox markazga joylashadi
    text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Matnning rasmning markaziga moslashtirilishi
    p = text_frame.add_paragraph()
    info = str(presentationEntity.Pages[pageNumber - 1].Text).replace("**", "").replace("\n\n", "\n")
    p.text = info
    size = 23
    lang = str(presentationEntity.Language).lower()
    if "rus" in lang:
        size = 20
    p.font.size = Pt(size)
    p.font.bold = False
    p.font.italic = False
    p.font.name = 'Vijaya'
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT  # Matn markazga moslashtirilishi

    if presentationEntity.Template in black_templates:
        # Slide ichidagi barcha shapeslarga murojat qilish
        for shape in slide1.shapes:
            if shape.has_text_frame:  # Faqat matn shakllariga murojat qiling
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)  # Oq rang