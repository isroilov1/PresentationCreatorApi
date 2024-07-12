from checkPhoto import open_image


async def add_page_slide_with_title(presentation, presentationEntity, black_templates, pageNumber):
    import os
    from PIL import Image

    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    # Page 3
    slide1_register = presentation.slide_layouts[0]
    slide1 = presentation.slides.add_slide(slide1_register)

    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = '...'
    subtitle1.text = 't.me/slaydai_bot'

    # Fon rasmini qo'shish
    bg_img_path = f"wwwroot/assets/templates/{presentationEntity.Template}/3.png"
    left_inch = Inches((presentation.slide_width - Inches(16)) / 2)
    top_inch = Inches((presentation.slide_height - Inches(9)) / 2)
    width_inch = Inches(16)
    height_inch = Inches(9)
    slide1.shapes.add_picture(bg_img_path, left_inch, top_inch, width_inch, height_inch)

    #  Plan title name
    text_box = slide1.shapes.add_textbox(Inches(1.5), Inches(0.7), Inches(7), Inches(2))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text_anchor = "middle"  # Matn textbox markazga joylashadi
    text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Matnning rasmning markaziga moslashtirilishi
    p = text_frame.add_paragraph()
    info = str(presentationEntity.Pages[pageNumber - 1].Title).replace("**", "").replace("\n\n", "\n")
    p.text = info
    p.font.size = Pt(27)
    p.font.bold = True
    p.font.italic = True
    p.font.name = 'ADLaM Display'
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT  # Matn markazga moslashtirilishi

    # Matn
    title_matn = presentationEntity.Pages[pageNumber - 1].Title
    top = 2.7
    if len(title_matn) < 60:
        top = 2.1
    elif len(title_matn) > 160:
        top += 1
    elif len(title_matn) > 130:
        top += 0.5
    text_box = slide1.shapes.add_textbox(Inches(1.5), Inches(top), Inches(7), Inches(4.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text_anchor = "top"  # Matn textbox markazga joylashadi
    text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Matnning rasmning markaziga moslashtirilishi
    p = text_frame.add_paragraph()
    info = str(presentationEntity.Pages[pageNumber - 1].Text).replace("**", "").replace("\n\n", "\n")
    p.text = info
    size = 23
    lang = str(presentationEntity.Language).lower()
    if "rus" in lang:
        size = 20
    p.font.size = Pt(size)
    p.font.bold = False
    p.font.italic = False
    p.font.name = 'Vijaya'
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT  # Matn markazga moslashtirilishi

    # Rasm yuklash
    image_num = 0
    if pageNumber < 4:
        image_num = 0
    elif pageNumber < 7:
        image_num = 1
    elif pageNumber < 10:
        image_num = 2
    elif pageNumber < 13:
        image_num = 3
    elif pageNumber < 16:
        image_num = 4
    elif pageNumber < 19:
        image_num = 5
    elif pageNumber < 22:
        image_num = 6
    image_path = f"wwwroot/{presentationEntity.ImagesPaths[image_num]}"
    if not os.path.exists(image_path):
        image_path = "wwwroot/assets/templates/1/1.png"

    with Image.open(image_path) as img:
        width, height = img.size
    top = 2
    width_img = 5
    left = 9.5
    if height / width > 3:
        top = 0.1
        width_img = 3
        left = 10.8
    elif height / width > 2:
        top = 0.3
        width_img = 4
        left = 9.5
    elif height / width > 1.5:
        top = 0.5
        width_img = 4.2
        left = 9.5
    elif height / width > 1:
        top = 1
        width_img = 4.5
        left = 9.5
    # Rasmni qo'shish
    slide1.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width_img))

    if presentationEntity.Template in black_templates:
        # Slide ichidagi barcha shapeslarga murojat qilish
        for shape in slide1.shapes:
            if shape.has_text_frame:  # Faqat matn shakllariga murojat qiling
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)  # Oq rang